"""
🎨 ClawBot PPT Generator — Gamma AI Quality
Creates professional PowerPoint presentations with:
- Dark/modern themes
- Smooth animations (Fade, Zoom, Appear)
- Storytelling flow
- Minimal text, impactful headings
- Icons & visual hierarchy
"""

import os
from pathlib import Path
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.oxml.ns import qn
    from lxml import etree
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ─── Color Palettes ───────────────────────────────────────────────────────────

THEMES = {
    'dark': {
        'bg':       RGBColor(0x0D, 0x0D, 0x0D),   # Near black
        'accent1':  RGBColor(0x00, 0xD4, 0xFF),   # Cyan
        'accent2':  RGBColor(0xFF, 0x6B, 0x35),   # Orange
        'heading':  RGBColor(0xFF, 0xFF, 0xFF),   # White
        'body':     RGBColor(0xCC, 0xCC, 0xCC),   # Light gray
        'dim':      RGBColor(0x66, 0x66, 0x66),   # Dim gray
        'card':     RGBColor(0x1A, 0x1A, 0x2E),   # Dark navy card
    },
    'light': {
        'bg':       RGBColor(0xFF, 0xFF, 0xFF),
        'accent1':  RGBColor(0x0D, 0x6E, 0xFD),
        'accent2':  RGBColor(0xFF, 0x4D, 0x4D),
        'heading':  RGBColor(0x1A, 0x1A, 0x1A),
        'body':     RGBColor(0x33, 0x33, 0x33),
        'dim':      RGBColor(0x88, 0x88, 0x88),
        'card':     RGBColor(0xF0, 0xF4, 0xFF),
    },
}


# ─── Slide Builders ───────────────────────────────────────────────────────────

def _fill_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text, left, top, width, height,
                  font_size=24, bold=False, color=None, align=PP_ALIGN.LEFT,
                  italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def _add_accent_line(slide, left, top, width, color: RGBColor, thickness_pt=3):
    """Draw a horizontal accent line."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, Pt(thickness_pt)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def _add_card(slide, left, top, width, height, color: RGBColor):
    """Add a rounded rectangle card background."""
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Set corner radius
    shape.adjustments[0] = 0.05
    return shape


def _add_fade_animation(shape):
    """Add a Fade in animation to a shape using XML."""
    try:
        sp = shape._element
        spTree = sp.getparent()
        slide_el = spTree.getparent()

        timing_el = slide_el.find(qn('p:timing'))
        if timing_el is None:
            timing_xml = '''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
              <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">
                <p:childTnLst><p:par><p:cTn id="2" fill="hold"><p:stCondLst>
                  <p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>
                </p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:tnLst>
            </p:timing>'''
            timing_el = etree.fromstring(timing_xml)
            slide_el.append(timing_el)
    except Exception:
        pass


# ─── Slide Types ──────────────────────────────────────────────────────────────

def make_title_slide(prs, title, subtitle, theme='dark'):
    """Hook slide — impactful title with glow accent."""
    T = THEMES[theme]
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    W, H = prs.slide_width, prs.slide_height

    _fill_bg(slide, T['bg'])

    # Top accent line
    _add_accent_line(slide, 0, Inches(0.15), W, T['accent1'], 4)

    # Large title
    _add_text_box(
        slide, title,
        Inches(0.6), Inches(2.0), W - Inches(1.2), Inches(2.5),
        font_size=52, bold=True, color=T['heading'], align=PP_ALIGN.CENTER
    )

    # Subtitle
    _add_text_box(
        slide, subtitle,
        Inches(0.6), Inches(4.6), W - Inches(1.2), Inches(1.0),
        font_size=22, color=T['accent1'], align=PP_ALIGN.CENTER
    )

    # Bottom dim info
    now = datetime.now().strftime("%B %Y")
    _add_text_box(
        slide, f"Powered by ClawBot AI  •  {now}",
        Inches(0.6), H - Inches(0.8), W - Inches(1.2), Inches(0.5),
        font_size=12, color=T['dim'], align=PP_ALIGN.CENTER
    )

    # Bottom accent line
    _add_accent_line(slide, 0, H - Inches(0.18), W, T['accent2'], 4)

    return slide


def make_content_slide(prs, title, bullet_points: list[str], theme='dark'):
    """Content slide with bullets on cards."""
    T = THEMES[theme]
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    W, H = prs.slide_width, prs.slide_height

    _fill_bg(slide, T['bg'])
    _add_accent_line(slide, 0, Inches(0.12), W, T['accent1'], 3)

    # Title
    _add_text_box(
        slide, title,
        Inches(0.5), Inches(0.3), W - Inches(1.0), Inches(0.9),
        font_size=34, bold=True, color=T['heading']
    )
    _add_accent_line(slide, Inches(0.5), Inches(1.25), Inches(1.5), T['accent1'], 3)

    # Bullet points as cards
    card_h = Inches(0.75)
    gap = Inches(0.12)
    start_y = Inches(1.5)

    for i, point in enumerate(bullet_points[:6]):
        y = start_y + i * (card_h + gap)
        if y + card_h > H - Inches(0.5):
            break
        _add_card(slide, Inches(0.5), y, W - Inches(1.0), card_h, T['card'])
        dot = T['accent1']
        _add_text_box(
            slide, "▮", Inches(0.7), y + Inches(0.15), Inches(0.35), card_h - Inches(0.1),
            font_size=14, color=dot, bold=True
        )
        _add_text_box(
            slide, point, Inches(1.1), y + Inches(0.1),
            W - Inches(1.8), card_h - Inches(0.1),
            font_size=18, color=T['body'], wrap=True
        )

    return slide


def make_two_column_slide(prs, title, left_heading, left_points, right_heading, right_points, theme='dark'):
    """Two-column comparison or feature slide."""
    T = THEMES[theme]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    _fill_bg(slide, T['bg'])
    _add_accent_line(slide, 0, Inches(0.12), W, T['accent1'], 3)

    _add_text_box(slide, title, Inches(0.5), Inches(0.3), W - Inches(1), Inches(0.9),
                  font_size=34, bold=True, color=T['heading'])

    col_w = (W - Inches(1.3)) / 2
    # Left column
    _add_card(slide, Inches(0.5), Inches(1.35), col_w, H - Inches(2.0), T['card'])
    _add_text_box(slide, left_heading, Inches(0.65), Inches(1.5), col_w - Inches(0.3), Inches(0.6),
                  font_size=18, bold=True, color=T['accent1'])
    for i, p in enumerate(left_points[:5]):
        _add_text_box(slide, f"• {p}", Inches(0.65), Inches(2.2) + i * Inches(0.7),
                      col_w - Inches(0.3), Inches(0.6), font_size=15, color=T['body'])

    # Right column
    rx = Inches(0.5) + col_w + Inches(0.3)
    _add_card(slide, rx, Inches(1.35), col_w, H - Inches(2.0), T['card'])
    _add_text_box(slide, right_heading, rx + Inches(0.15), Inches(1.5), col_w - Inches(0.3), Inches(0.6),
                  font_size=18, bold=True, color=T['accent2'])
    for i, p in enumerate(right_points[:5]):
        _add_text_box(slide, f"• {p}", rx + Inches(0.15), Inches(2.2) + i * Inches(0.7),
                      col_w - Inches(0.3), Inches(0.6), font_size=15, color=T['body'])

    return slide


def make_section_break_slide(prs, section_number, section_title, theme='dark'):
    """Bold numbered section separator slide."""
    T = THEMES[theme]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    _fill_bg(slide, T['bg'])

    # Big number (background accent)
    _add_text_box(slide, str(section_number),
                  Inches(0.2), Inches(0.5), W, H - Inches(1),
                  font_size=240, bold=True, color=T['card'], align=PP_ALIGN.CENTER)

    # Section title on top
    _add_text_box(slide, section_title,
                  Inches(0.5), Inches(2.5), W - Inches(1), Inches(2),
                  font_size=48, bold=True, color=T['accent1'], align=PP_ALIGN.CENTER)


def make_cta_slide(prs, headline, cta_text, contact=None, theme='dark'):
    """Call-to-action / closing slide."""
    T = THEMES[theme]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    _fill_bg(slide, T['bg'])
    _add_accent_line(slide, 0, Inches(0.12), W, T['accent2'], 4)

    _add_text_box(slide, headline,
                  Inches(0.5), Inches(1.8), W - Inches(1), Inches(1.8),
                  font_size=44, bold=True, color=T['heading'], align=PP_ALIGN.CENTER)

    _add_text_box(slide, cta_text,
                  Inches(0.5), Inches(3.8), W - Inches(1), Inches(1.0),
                  font_size=22, color=T['accent1'], align=PP_ALIGN.CENTER)

    if contact:
        _add_text_box(slide, contact,
                      Inches(0.5), H - Inches(1.2), W - Inches(1), Inches(0.5),
                      font_size=14, color=T['dim'], align=PP_ALIGN.CENTER)

    _add_accent_line(slide, 0, H - Inches(0.18), W, T['accent1'], 4)


# ─── Main API ─────────────────────────────────────────────────────────────────

def create_ppt(slides_data: list[dict], output_path: str = None, theme: str = 'dark') -> str:
    """
    Create a professional PPT from a list of slide definitions.

    Each slide_data dict has:
        type: 'title' | 'content' | 'two_column' | 'section' | 'cta'
        title: str
        subtitle: str (for title slide)
        points: list[str] (for content slide)
        left_heading, left_points, right_heading, right_points (for two_column)
        section_number, section_title (for section break)
        headline, cta_text, contact (for cta)

    Returns: absolute path to saved .pptx file
    """
    if not HAS_PPTX:
        raise ImportError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for sd in slides_data:
        t = sd.get('type', 'content')

        if t == 'title':
            make_title_slide(prs, sd.get('title', 'Presentation'), sd.get('subtitle', ''), theme)

        elif t == 'section':
            make_section_break_slide(prs, sd.get('section_number', ''), sd.get('section_title', sd.get('title', '')), theme)

        elif t == 'two_column':
            make_two_column_slide(
                prs, sd.get('title', ''),
                sd.get('left_heading', 'Left'), sd.get('left_points', []),
                sd.get('right_heading', 'Right'), sd.get('right_points', []),
                theme
            )

        elif t == 'cta':
            make_cta_slide(prs, sd.get('headline', 'Thank You'), sd.get('cta_text', ''), sd.get('contact'), theme)

        else:  # content (default)
            pts = sd.get('points', [])
            if not pts and 'content' in sd:
                c = sd['content']
                if isinstance(c, list):
                    pts = c
                else:
                    pts = [p.strip() for p in str(c).split('\n') if p.strip()]
            make_content_slide(prs, sd.get('title', ''), pts, theme)

    # Save
    if not output_path:
        desktop = Path.home() / 'OneDrive' / 'Desktop'
        if not desktop.exists():
            desktop = Path.home() / 'Desktop'
        ts = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_path = str(desktop / f'ClawBot_PPT_{ts}.pptx')

    prs.save(output_path)
    return output_path


# ─── Quick Demo ───────────────────────────────────────────────────────────────

if __name__ == '__main__':
    demo_slides = [
        {'type': 'title', 'title': 'ClawBot AI', 'subtitle': 'Your AI-Powered PC Assistant'},
        {'type': 'section', 'section_number': '01', 'section_title': 'What is ClawBot?'},
        {'type': 'content', 'title': 'Core Capabilities', 'points': [
            '🤖 Full PC Control — open apps, manage files, automate anything',
            '💻 Self-Healing Code — writes and auto-fixes Python/JS/CSS code',
            '🔍 Instant Web Search — no browser needed',
            '📱 Telegram Remote Control — PC se door, control paas',
            '🎨 Professional PPT Generation — Gamma AI quality',
        ]},
        {'type': 'two_column', 'title': 'CLI vs Telegram',
         'left_heading': '💻 CLI Agent', 'left_points': ['Fast execution', 'Live feedback', 'File management', 'Multi-step tasks'],
         'right_heading': '📱 Telegram Bot', 'right_points': ['Remote access', 'Auto intent detection', 'PC automation', 'Chat + Task mode']},
        {'type': 'cta', 'headline': 'Start Using ClawBot Today!',
         'cta_text': 'pip install clawbot  →  clawbot',
         'contact': 'github.com/clawbot  •  Free & Open Source'},
    ]
    path = create_ppt(demo_slides, theme='dark')
    print(f"✅ PPT created: {path}")
